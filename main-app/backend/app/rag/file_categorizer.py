"""
FileCategorizationService — Phase 2 of the revamped training pipeline.

Determines the content category of each file using only local inspection
(no LLM calls).  Categories:

  text           — content is directly extractable as UTF-8 text
  image_only     — content is purely image-based (no extractable text)
  text_in_image  — content *looks* like images but contains OCR-extractable text
  mixed          — file contains both extractable text AND embedded images
  exclude        — not text, not image, not useful (dropped from pipeline)
"""

from typing import Optional

import structlog

from app.rag.models import FileRecord, CategorizedFile, ContentCategory

logger = structlog.get_logger()

# ── Extension sets ────────────────────────────────────────────────────────

PURE_TEXT_EXTENSIONS: set[str] = {
    # Source code
    ".py", ".pyi", ".js", ".jsx", ".ts", ".tsx", ".java", ".go", ".rb",
    ".cs", ".cpp", ".c", ".h", ".hpp", ".rs", ".swift", ".kt", ".kts",
    ".scala", ".php", ".sh", ".bash", ".zsh", ".sql", ".r", ".lua",
    ".ex", ".exs", ".clj", ".hs", ".dart", ".m", ".mm", ".pl", ".pm",
    ".ps1", ".bat", ".cmd", ".fish", ".groovy", ".v", ".vhdl",
    # Markup / docs
    ".md", ".mdx", ".rst", ".txt", ".html", ".htm", ".xml", ".xhtml",
    ".tex", ".adoc", ".asciidoc", ".org", ".wiki", ".textile",
    # Data / config
    ".json", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf", ".env",
    ".env.example", ".env.sample", ".properties", ".csv", ".tsv",
    ".editorconfig", ".prettierrc", ".eslintrc", ".babelrc",
    # Build / CI
    ".dockerfile", ".makefile", ".cmake", ".gradle",
    ".tf", ".hcl", ".bicep",
    # Other text
    ".rtf", ".log", ".diff", ".patch",
}

PURE_IMAGE_EXTENSIONS: set[str] = {
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".svg",
    ".drawio", ".ico", ".tiff", ".tif",
}

DOCUMENT_EXTENSIONS: set[str] = {
    ".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls",
    ".vsdx", ".vsd", ".odt", ".ods", ".odp", ".pages", ".numbers",
    ".key", ".epub",
}


class FileCategorizationService:
    """Categorize files by content type using local inspection only."""

    def __init__(self):
        self._stats = {
            "text": 0, "image_only": 0, "text_in_image": 0,
            "mixed": 0, "exclude": 0, "errors": 0,
        }

    @property
    def stats(self) -> dict[str, int]:
        return dict(self._stats)

    def categorize(
        self,
        files: list[FileRecord],
        on_file: Optional[callable] = None,
    ) -> list[CategorizedFile]:
        """Categorize a list of files.  Returns CategorizedFile list (including excludes).

        *on_file* is called with (index, file, category_str) for progress.
        """
        self._stats = {
            "text": 0, "image_only": 0, "text_in_image": 0,
            "mixed": 0, "exclude": 0, "errors": 0,
        }
        results: list[CategorizedFile] = []

        for idx, f in enumerate(files):
            try:
                cat = self._categorize_one(f)
            except Exception as exc:
                logger.warning("categorize_error", path=f.path, error=str(exc))
                cat = CategorizedFile(
                    **f.model_dump(),
                    content_category=ContentCategory.EXCLUDE,
                )
                self._stats["errors"] += 1

            self._stats[cat.content_category.value] += 1
            results.append(cat)

            if on_file:
                try:
                    on_file(idx, f, cat.content_category.value)
                except Exception:
                    pass

        logger.info("categorization_complete", total=len(results), **self._stats)
        return results

    def _categorize_one(self, f: FileRecord) -> CategorizedFile:
        ext = f.ext.lower()
        base = f.model_dump()

        if ext in PURE_TEXT_EXTENSIONS:
            return CategorizedFile(**base, content_category=ContentCategory.TEXT)

        if ext in PURE_IMAGE_EXTENSIONS:
            return CategorizedFile(**base, content_category=ContentCategory.IMAGE_ONLY)

        if ext in DOCUMENT_EXTENSIONS:
            return self._probe_document(f)

        # Unknown extension — try reading as UTF-8
        return self._probe_unknown(f)

    def _probe_document(self, f: FileRecord) -> CategorizedFile:
        """Probe a document file (PDF, DOCX, etc.) to determine content type."""
        ext = f.ext.lower()
        base = f.model_dump()

        if ext == ".pdf":
            return self._probe_pdf(f)
        elif ext in (".docx", ".doc"):
            return self._probe_docx(f)
        elif ext in (".pptx", ".ppt"):
            return self._probe_pptx(f)
        elif ext in (".xlsx", ".xls"):
            return CategorizedFile(**base, content_category=ContentCategory.TEXT)
        elif ext in (".vsdx", ".vsd"):
            return CategorizedFile(**base, content_category=ContentCategory.IMAGE_ONLY)
        elif ext == ".epub":
            return CategorizedFile(**base, content_category=ContentCategory.TEXT)
        else:
            return self._probe_unknown(f)

    def _probe_pdf(self, f: FileRecord) -> CategorizedFile:
        """Probe a PDF: try text extraction and image detection. Prefer PyMuPDF (faster)."""
        base = f.model_dump()

        text_length = 0
        has_images = False
        page_count = 0

        try:
            import fitz  # pymupdf
        except ImportError:
            fitz = None
        if fitz is not None:
            try:
                doc = fitz.open(f.path)
                try:
                    page_count = len(doc)
                    for page in doc:
                        text_length += len((page.get_text() or ""))
                        if page.get_images():
                            has_images = True
                finally:
                    doc.close()
                return self._probe_pdf_result(base, text_length, has_images, page_count, f.path)
            except Exception as exc:
                logger.debug("pdf_probe_failed", path=f.path, error=str(exc), backend="pymupdf")

        try:
            import pdfplumber
            with pdfplumber.open(f.path) as pdf:
                page_count = len(pdf.pages)
                for page in pdf.pages:
                    page_text = page.extract_text() or ""
                    text_length += len(page_text)
                    if page.images:
                        has_images = True
        except Exception as exc:
            logger.debug("pdf_probe_failed", path=f.path, error=str(exc))
            return CategorizedFile(**base, content_category=ContentCategory.EXCLUDE)

        return self._probe_pdf_result(base, text_length, has_images, page_count, f.path)

    def _probe_pdf_result(
        self, base: dict, text_length: int, has_images: bool, page_count: int, path: str
    ) -> CategorizedFile:

        has_text = text_length > 50

        if has_text and has_images:
            return CategorizedFile(
                **base,
                content_category=ContentCategory.MIXED,
                estimated_text_length=text_length,
            )
        elif has_text:
            return CategorizedFile(
                **base,
                content_category=ContentCategory.TEXT,
                estimated_text_length=text_length,
            )
        elif has_images:
            ocr_pass = self._ocr_probe_pdf(path, page_count)
            if ocr_pass:
                return CategorizedFile(
                    **base,
                    content_category=ContentCategory.TEXT_IN_IMAGE,
                    ocr_required=True,
                )
            return CategorizedFile(**base, content_category=ContentCategory.IMAGE_ONLY)
        else:
            return CategorizedFile(**base, content_category=ContentCategory.EXCLUDE)

    def _ocr_probe_pdf(self, path: str, page_count: int) -> bool:
        """Sample 2 random pages, run OCR, check if >= 100 words extracted."""
        if page_count == 0:
            return False
        try:
            from app.rag.ocr_service import ocr_probe_pdf
            return ocr_probe_pdf(path, page_count)
        except ImportError:
            logger.debug("ocr_service_not_available")
            return False
        except Exception as exc:
            logger.debug("ocr_probe_failed", path=path, error=str(exc))
            return False

    def _probe_docx(self, f: FileRecord) -> CategorizedFile:
        base = f.model_dump()
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(f.path)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            has_images = any(
                rel.reltype.endswith("/image")
                for rel in doc.part.rels.values()
            )
            text_length = len(text)

            if text_length > 50 and has_images:
                return CategorizedFile(**base, content_category=ContentCategory.MIXED, estimated_text_length=text_length)
            elif text_length > 50:
                return CategorizedFile(**base, content_category=ContentCategory.TEXT, estimated_text_length=text_length)
            elif has_images:
                return CategorizedFile(**base, content_category=ContentCategory.IMAGE_ONLY)
            else:
                return CategorizedFile(**base, content_category=ContentCategory.EXCLUDE)
        except Exception as exc:
            logger.debug("docx_probe_failed", path=f.path, error=str(exc))
            return CategorizedFile(**base, content_category=ContentCategory.EXCLUDE)

    def _probe_pptx(self, f: FileRecord) -> CategorizedFile:
        base = f.model_dump()
        try:
            from pptx import Presentation
            prs = Presentation(f.path)
            text_length = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_length += len(shape.text)

            if text_length > 50:
                return CategorizedFile(**base, content_category=ContentCategory.MIXED, estimated_text_length=text_length)
            else:
                return CategorizedFile(**base, content_category=ContentCategory.IMAGE_ONLY)
        except ImportError:
            return CategorizedFile(**base, content_category=ContentCategory.EXCLUDE)
        except Exception:
            return CategorizedFile(**base, content_category=ContentCategory.EXCLUDE)

    def _probe_unknown(self, f: FileRecord) -> CategorizedFile:
        """Try reading as UTF-8 text; if decodable → text, else → exclude."""
        base = f.model_dump()
        try:
            with open(f.path, "r", encoding="utf-8") as fh:
                sample = fh.read(4096)
            printable_ratio = sum(1 for c in sample if c.isprintable() or c in "\n\r\t") / max(len(sample), 1)
            if printable_ratio > 0.85 and len(sample.strip()) > 10:
                return CategorizedFile(
                    **base,
                    content_category=ContentCategory.TEXT,
                    estimated_text_length=f.size_bytes,
                )
            return CategorizedFile(**base, content_category=ContentCategory.EXCLUDE)
        except (UnicodeDecodeError, OSError):
            return CategorizedFile(**base, content_category=ContentCategory.EXCLUDE)

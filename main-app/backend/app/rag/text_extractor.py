"""
TextExtractorService — Phase 3: extract raw text from kept files.

Dispatch by content_type assigned by the LLM in Phase 2:
  text           → direct UTF-8 read (or specialised parser for html, etc.)
  text_or_image  → try direct extraction first; if quality is poor, OCR probe
                   then full OCR
  image          → skip (no text extraction; images kept for future vision use)

Each extraction method is logged with detail so the user knows exactly what
is happening for each file.
"""

import re
from typing import Optional

import structlog

from app.rag.models import TreeNode

logger = structlog.get_logger()


def _human_size(nbytes: int) -> str:
    for unit in ("B", "KB", "MB", "GB"):
        if abs(nbytes) < 1024:
            return f"{nbytes:.1f}{unit}" if unit != "B" else f"{nbytes}{unit}"
        nbytes /= 1024  # type: ignore[assignment]
    return f"{nbytes:.1f}TB"


def _count_real_words(text: str) -> int:
    """Count words that look like natural language (>= 2 alpha chars)."""
    return len(re.findall(r"\b[a-zA-Z]{2,}\b", text))


class TextExtractorService:
    """Extract raw text from files based on their LLM-assigned content_type."""

    def __init__(self):
        self.stats = {
            "text_direct": 0,
            "text_or_image_direct": 0,
            "text_or_image_ocr": 0,
            "text_or_image_failed": 0,
            "image_skipped": 0,
            "errors": 0,
            "total_chars": 0,
        }

    def extract(self, node: TreeNode, on_progress=None) -> str:
        """Extract text from a single TreeNode file.

        Mutates node.extracted_text, node.extraction_method, node.extraction_chars.
        Returns the extracted text.
        """
        if node.decision != "keep" or not node.is_file:
            return ""

        ct = node.content_type
        try:
            if ct == "text":
                text = self._extract_text(node)
                node.extraction_method = "direct"
            elif ct == "text_or_image":
                text = self._extract_mixed(node, on_progress)
            elif ct == "image":
                self.stats["image_skipped"] += 1
                node.extraction_method = "skipped"
                return ""
            else:
                text = self._extract_text(node)
                node.extraction_method = "direct"

            node.extracted_text = text
            node.extraction_chars = len(text)
            if text.strip():
                self.stats["total_chars"] += len(text)
            return text

        except Exception as exc:
            logger.error("text_extract_error", path=node.path, error=str(exc))
            self.stats["errors"] += 1
            node.extraction_method = "failed"
            return ""

    # ------------------------------------------------------------------
    # Text files — direct read
    # ------------------------------------------------------------------

    def _extract_text(self, node: TreeNode) -> str:
        """Read a definitely-text file."""
        ext = node.ext.lower()

        if ext == ".pdf":
            text = self._extract_pdf_text(node.path)
        elif ext in (".docx", ".doc"):
            text = self._extract_docx(node.path)
        elif ext in (".html", ".htm", ".xhtml"):
            text = self._extract_html(node.path)
        elif ext in (".xlsx", ".xls"):
            text = self._extract_xlsx(node.path)
        elif ext in (".pptx", ".ppt"):
            text = self._extract_pptx(node.path)
        elif ext == ".epub":
            text = self._extract_epub(node.path)
        elif ext == ".ipynb":
            text = self._extract_ipynb(node.path)
        elif ext == ".rtf":
            text = self._extract_rtf(node.path)
        else:
            text = self._read_utf8(node.path)

        self.stats["text_direct"] += 1
        return text

    # ------------------------------------------------------------------
    # Mixed files — direct first, OCR fallback
    # ------------------------------------------------------------------

    def _extract_mixed(self, node: TreeNode, on_progress=None) -> str:
        """Try direct extraction; if too little text, fall back to OCR."""
        # Step 1: direct extraction
        text = self._extract_text(node)
        self.stats["text_direct"] -= 1  # undo count from _extract_text

        if text.strip() and len(text.strip()) > 200 and _count_real_words(text) > 30:
            self.stats["text_or_image_direct"] += 1
            node.extraction_method = "direct"
            logger.info(
                "mixed_direct_ok",
                path=node.path,
                chars=len(text),
                words=_count_real_words(text),
            )
            return text

        # Step 2: OCR probe (for PDFs)
        if node.ext.lower() == ".pdf":
            return self._ocr_pdf_with_probe(node, on_progress)

        # For non-PDF mixed files, direct extraction is our only option
        if text.strip():
            self.stats["text_or_image_direct"] += 1
            node.extraction_method = "direct"
            return text

        self.stats["text_or_image_failed"] += 1
        node.extraction_method = "failed"
        return ""

    def _ocr_pdf_with_probe(self, node: TreeNode, on_progress=None) -> str:
        """OCR probe 1-2 pages; if promising, OCR the full document."""
        from app.rag.ocr_service import ocr_probe_pdf, ocr_full_pdf, _check_tesseract

        if not _check_tesseract():
            logger.info("ocr_not_available", path=node.path)
            self.stats["text_or_image_failed"] += 1
            node.extraction_method = "failed"
            return ""

        # Get page count for the probe
        page_count = self._pdf_page_count(node.path)
        if page_count <= 0:
            self.stats["text_or_image_failed"] += 1
            node.extraction_method = "failed"
            return ""

        logger.info("ocr_probe_start", path=node.path, pages=page_count)
        has_text = ocr_probe_pdf(node.path, page_count)

        if not has_text:
            logger.info("ocr_probe_negative", path=node.path)
            self.stats["text_or_image_failed"] += 1
            node.extraction_method = "failed"
            return ""

        # Probe positive — full OCR
        logger.info("ocr_full_start", path=node.path, pages=page_count)

        def _on_page(page_num, total):
            if on_progress:
                try:
                    on_progress(
                        f"OCR page {page_num}/{total} of {node.name}…"
                    )
                except Exception:
                    pass

        text = ocr_full_pdf(node.path, on_page=_on_page)
        if text.strip():
            self.stats["text_or_image_ocr"] += 1
            node.extraction_method = "ocr"
            logger.info("ocr_full_done", path=node.path, chars=len(text))
            return text

        self.stats["text_or_image_failed"] += 1
        node.extraction_method = "failed"
        return ""

    # ------------------------------------------------------------------
    # Format-specific extractors
    # ------------------------------------------------------------------

    def _read_utf8(self, path: str, max_bytes: int = 1024 * 1024 * 1024) -> str:
        try:
            with open(path, "r", encoding="utf-8", errors="replace") as fh:
                return fh.read(max_bytes)
        except Exception:
            return ""

    def _extract_pdf_text(self, path: str) -> str:
        # Prefer PyMuPDF (faster); fall back to pdfplumber
        try:
            import fitz  # pymupdf
        except ImportError:
            fitz = None
        if fitz is not None:
            try:
                pages: list[str] = []
                doc = fitz.open(path)
                try:
                    for page in doc:
                        text = page.get_text() or ""
                        if text.strip():
                            pages.append(text)
                finally:
                    doc.close()
                return "\n\n".join(pages)
            except Exception as exc:
                logger.warning("pdf_extract_error", path=path, error=str(exc), backend="pymupdf")
                return ""
        try:
            import pdfplumber
        except ImportError:
            logger.warning("pdfplumber_not_installed")
            return ""
        pages = []
        try:
            with pdfplumber.open(path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text() or ""
                    if text.strip():
                        pages.append(text)
        except Exception as exc:
            logger.warning("pdf_extract_error", path=path, error=str(exc))
        return "\n\n".join(pages)

    def _extract_docx(self, path: str) -> str:
        try:
            from docx import Document as DocxDocument
        except ImportError:
            logger.warning("python_docx_not_installed")
            return ""
        try:
            doc = DocxDocument(path)
            return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception as exc:
            logger.warning("docx_extract_error", path=path, error=str(exc))
            return ""

    def _extract_html(self, path: str) -> str:
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            return self._read_utf8(path)
        try:
            raw = self._read_utf8(path)
            soup = BeautifulSoup(raw, "html.parser")
            return soup.get_text(separator="\n")
        except Exception:
            return self._read_utf8(path)

    def _extract_xlsx(self, path: str) -> str:
        try:
            import openpyxl
        except ImportError:
            return ""
        try:
            wb = openpyxl.load_workbook(path, read_only=True)
            rows: list[str] = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    vals = [str(v) for v in row if v is not None]
                    if vals:
                        rows.append("\t".join(vals))
            wb.close()
            return "\n".join(rows)
        except Exception as exc:
            logger.warning("xlsx_extract_error", path=path, error=str(exc))
            return ""

    def _extract_pptx(self, path: str) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            return ""
        try:
            prs = Presentation(path)
            parts: list[str] = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text: list[str] = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))
            return "\n\n".join(parts)
        except Exception:
            return ""

    def _extract_epub(self, path: str) -> str:
        try:
            import zipfile
            from bs4 import BeautifulSoup
            parts: list[str] = []
            with zipfile.ZipFile(path, "r") as zf:
                for name in zf.namelist():
                    if name.endswith((".html", ".xhtml", ".htm")):
                        raw = zf.read(name).decode("utf-8", errors="replace")
                        soup = BeautifulSoup(raw, "html.parser")
                        text = soup.get_text(separator="\n").strip()
                        if text:
                            parts.append(text)
            return "\n\n".join(parts)
        except Exception:
            return ""

    def _extract_ipynb(self, path: str) -> str:
        """Extract text from Jupyter notebooks."""
        import json as _json
        try:
            with open(path, "r", encoding="utf-8") as fh:
                nb = _json.load(fh)
            cells = nb.get("cells", [])
            parts: list[str] = []
            for cell in cells:
                cell_type = cell.get("cell_type", "")
                source = "".join(cell.get("source", []))
                if source.strip():
                    parts.append(f"[{cell_type}]\n{source}")
            return "\n\n".join(parts)
        except Exception:
            return ""

    def _extract_rtf(self, path: str) -> str:
        try:
            from striprtf.striprtf import rtf_to_text
            raw = self._read_utf8(path)
            return rtf_to_text(raw)
        except ImportError:
            return self._read_utf8(path)
        except Exception:
            return ""

    def _pdf_page_count(self, path: str) -> int:
        try:
            import fitz
            doc = fitz.open(path)
            try:
                return len(doc)
            finally:
                doc.close()
        except ImportError:
            pass
        except Exception:
            pass
        try:
            import pdfplumber
            with pdfplumber.open(path) as pdf:
                return len(pdf.pages)
        except Exception:
            return 0
